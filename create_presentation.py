"""Cria apresentação PowerPoint sobre o sistema de Adaptive HTTP Streaming."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de cores ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # azul-marinho escuro
BG_CARD      = RGBColor(0x16, 0x2A, 0x3F)   # card ligeiramente mais claro
ACCENT_BLUE  = RGBColor(0x1F, 0x77, 0xB4)   # azul primário (mesma do gráfico)
ACCENT_GREEN = RGBColor(0x2C, 0xA0, 0x2C)   # verde (quality timeline)
ACCENT_PURP  = RGBColor(0x94, 0x67, 0xBD)   # roxo (buffer level)
ACCENT_ORANGE= RGBColor(0xFD, 0xAE, 0x6B)   # laranja (histograma)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)
CODE_BG      = RGBColor(0x0A, 0x12, 0x1C)   # fundo de bloco de código
CODE_TEXT    = RGBColor(0x7E, 0xC8, 0xE3)   # texto de código azul-claro

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────
def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=20, bold=False,
             color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, title, bullets, x, y, w, h,
                   title_color=ACCENT_BLUE, bullet_color=LIGHT_GRAY,
                   title_size=18, bullet_size=15, bg_color=None):
    if bg_color:
        add_rect(slide, x, y, w, h, bg_color)
    add_text(slide, title, x + 0.15, y + 0.1, w - 0.3, 0.4,
             size=title_size, bold=True, color=title_color)
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.55), Inches(w - 0.3), Inches(h - 0.65)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = bullet_color
        p.space_after = Pt(4)


def add_code_block(slide, code_lines, x, y, w, h, font_size=10.5):
    add_rect(slide, x, y, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(h - 0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = CODE_TEXT
        run.font.name = "Courier New"


def add_accent_bar(slide, color=ACCENT_BLUE):
    """Barra de cor no topo do slide."""
    add_rect(slide, 0, 0, 13.33, 0.08, color)


def slide_number_label(slide, n):
    add_text(slide, str(n), 12.7, 7.1, 0.5, 0.3,
             size=11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ── Criação dos slides ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # completamente em branco


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Capa
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)

# Gradiente simulado com retângulos
add_rect(sl, 0, 0, 13.33, 7.5, RGBColor(0x08, 0x12, 0x20))
add_rect(sl, 0, 5.8, 13.33, 1.7, RGBColor(0x12, 0x22, 0x38))

# Linha decorativa
add_rect(sl, 1.0, 2.3, 11.33, 0.05, ACCENT_BLUE)

# Ícone / badge
add_rect(sl, 5.9, 0.8, 1.53, 1.35, ACCENT_BLUE)
add_text(sl, "TR2", 5.9, 0.85, 1.53, 1.1,
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Título principal
add_text(sl, "Adaptive HTTP Streaming",
         1.0, 2.5, 11.33, 1.2,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtítulo
add_text(sl, "Streaming Adaptativo com ABR — Arquitetura, Algoritmos e Análise",
         1.0, 3.75, 11.33, 0.7,
         size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)

# Linha decorativa inferior
add_rect(sl, 1.0, 4.55, 11.33, 0.05, ACCENT_GREEN)

# Rodapé
add_text(sl, "Teleinformática e Redes 2  |  Projeto Final",
         1.0, 6.8, 11.33, 0.5,
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Agenda
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_BLUE)
slide_number_label(sl, 2)

add_text(sl, "Agenda", 0.4, 0.2, 12.0, 0.7,
         size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 4.5, 0.04, ACCENT_BLUE)

topics = [
    ("1", "O que é Adaptive HTTP Streaming?",    ACCENT_BLUE),
    ("2", "Arquitetura do Sistema",               ACCENT_GREEN),
    ("3", "ManifestParser — Negociação de Qualidades", ACCENT_PURP),
    ("4", "RateBasedABR — Algoritmo de Seleção",  ACCENT_BLUE),
    ("5", "BufferManager — Gestão do Buffer",      ACCENT_ORANGE),
    ("6", "ThroughputMeter — Medição de Banda",    ACCENT_GREEN),
    ("7", "MetricsRecorder — Coleta de Dados",     ACCENT_PURP),
    ("8", "run_simulation.py — Simulação Completa",ACCENT_ORANGE),
    ("9", "Visualização: Gráficos e Resultados",   ACCENT_BLUE),
]

for i, (num, title, color) in enumerate(topics):
    row = i % 5
    col = i // 5
    bx = 0.5 + col * 6.4
    by = 1.3 + row * 1.1

    add_rect(sl, bx, by, 6.1, 0.9, BG_CARD)
    add_rect(sl, bx, by, 0.55, 0.9, color)
    add_text(sl, num, bx, by, 0.55, 0.9,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, bx + 0.65, by + 0.15, 5.3, 0.6,
             size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – O que é ABR / Contexto
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_BLUE)
slide_number_label(sl, 3)

add_text(sl, "O que é Adaptive HTTP Streaming?",
         0.4, 0.2, 12.0, 0.7, size=28, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.5, 0.04, ACCENT_BLUE)

# Coluna esquerda — definição
add_bullet_box(
    sl,
    "Conceito",
    [
        "Vídeo dividido em segmentos de duração fixa (ex.: 4 s)",
        "Cada segmento disponível em múltiplas qualidades (bitrates)",
        "Cliente escolhe a qualidade em tempo real com base na rede",
        "Protocolo de transporte: HTTP sobre TCP",
        "Padrão de mercado: MPEG-DASH e HLS",
    ],
    0.4, 1.1, 5.8, 3.5,
    title_color=ACCENT_BLUE, bg_color=BG_CARD
)

# Coluna direita — problema
add_bullet_box(
    sl,
    "Por que é necessário?",
    [
        "Banda da internet varia constantemente",
        "Streaming estático causa rebuffering ou baixa qualidade fixa",
        "ABR busca maximizar qualidade sem pausas",
        "Métricas chave: throughput, buffer, switches de qualidade",
    ],
    6.4, 1.1, 6.5, 2.7,
    title_color=ACCENT_GREEN, bg_color=BG_CARD
)

# Diagrama simplificado de fluxo
add_rect(sl, 6.4, 4.0, 6.5, 3.0, BG_CARD)
add_text(sl, "Fluxo Básico", 6.55, 4.1, 6.2, 0.4,
         size=14, bold=True, color=ACCENT_ORANGE)

boxes = [
    (6.55, 4.65, "Servidor", ACCENT_BLUE),
    (8.35, 4.65, "Manifesto", ACCENT_PURP),
    (10.15, 4.65, "Segmento", ACCENT_GREEN),
    (11.9,  4.65, "Player",   ACCENT_ORANGE),
]
for bx, by, label, c in boxes:
    add_rect(sl, bx, by, 1.55, 0.55, c)
    add_text(sl, label, bx, by, 1.55, 0.55,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Setas
for bx in [8.12, 9.92, 11.67]:
    add_text(sl, "→", bx, 4.62, 0.25, 0.55,
             size=18, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(sl,
         "ABR decide qual Segmento baixar a cada ciclo\n"
         "com base no throughput medido e no buffer atual.",
         6.55, 5.45, 6.2, 1.3,
         size=13, color=LIGHT_GRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Arquitetura
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_GREEN)
slide_number_label(sl, 4)

add_text(sl, "Arquitetura do Sistema",
         0.4, 0.2, 12.0, 0.7, size=28, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 4.0, 0.04, ACCENT_GREEN)

modules = [
    ("manifest_parser.py", "ManifestParser",   "Baixa e interpreta o\nmanifesto JSON do servidor",      ACCENT_BLUE,   0.5,  1.3),
    ("abr.py",             "RateBasedABR",     "Seleciona a qualidade\ncom base no throughput",           ACCENT_GREEN,  4.1,  1.3),
    ("buffer_manager.py",  "BufferManager",    "Adiciona/consome segmentos\ne detecta rebuffering",       ACCENT_PURP,   7.7,  1.3),
    ("metrics.py",         "ThroughputMeter",  "Mede a banda de download\ncada segmento",                 ACCENT_ORANGE, 0.5,  4.2),
    ("metrics.py",         "MetricsRecorder",  "Persiste métricas em CSV\npor lote (batch)",              ACCENT_BLUE,   4.1,  4.2),
    ("graphs.py",          "generate_graphs",  "Lê CSV e gera 5 gráficos\nem PNG via Matplotlib",         ACCENT_GREEN,  7.7,  4.2),
]

for fname, cls, desc, color, bx, by in modules:
    add_rect(sl, bx, by, 5.4, 2.5, BG_CARD)
    add_rect(sl, bx, by, 5.4, 0.38, color)
    add_text(sl, fname, bx + 0.12, by + 0.05, 5.1, 0.28,
             size=11, color=WHITE, italic=True)
    add_text(sl, cls, bx + 0.12, by + 0.5, 5.1, 0.55,
             size=18, bold=True, color=WHITE)
    add_text(sl, desc, bx + 0.12, by + 1.1, 5.1, 1.2,
             size=14, color=LIGHT_GRAY)

# Conector central
add_text(sl, "run_simulation.py  |  player.py",
         3.5, 3.65, 6.33, 0.45,
         size=16, bold=True, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
add_rect(sl, 3.5, 3.6, 6.33, 0.04, ACCENT_ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – ManifestParser
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_PURP)
slide_number_label(sl, 5)

add_text(sl, "ManifestParser — Negociação de Qualidades",
         0.4, 0.2, 12.0, 0.7, size=26, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.5, 0.04, ACCENT_PURP)

add_bullet_box(
    sl,
    "Responsabilidades",
    [
        "download_manifest(url) — GET HTTP com timeout de 10 s",
        "load_from_dict(data) — carrega manifesto para testes",
        "get_qualities() — lista de {name, bitrate} ordenada",
        "get_servers() — servidores ordenados por prioridade",
        "get_segment_duration() — duração de cada chunk",
        "get_primary_server() — servidor de maior prioridade",
    ],
    0.4, 1.1, 5.8, 3.8,
    title_color=ACCENT_PURP, bg_color=BG_CARD
)

add_code_block(sl, [
    "# Estrutura esperada do manifesto JSON",
    "{",
    '  "qualities": [',
    '    {"name": "240p",  "bitrate": 200},',
    '    {"name": "720p",  "bitrate": 1200},',
    '    {"name": "1080p", "bitrate": 2500}',
    "  ],",
    '  "servers": [',
    '    {"url": "http://cdn1.example.com", "priority": 1},',
    '    {"url": "http://cdn2.example.com", "priority": 2}',
    "  ],",
    '  "segment": {"duration": 4.0, "total": 30}',
    "}",
], 6.4, 1.1, 6.5, 4.5, font_size=10.5)

add_bullet_box(
    sl,
    "Tratamento de Erros",
    [
        "ConnectionError → servidor indisponível",
        "Timeout → requisição expirada (10 s)",
        "HTTPError → status 4xx / 5xx",
        "JSONDecodeError → manifesto malformado",
        "KeyError → campo obrigatório ausente",
    ],
    0.4, 5.1, 5.8, 2.2,
    title_color=ACCENT_ORANGE, bg_color=BG_CARD, bullet_size=13
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – RateBasedABR
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_BLUE)
slide_number_label(sl, 6)

add_text(sl, "RateBasedABR — Algoritmo de Seleção de Qualidade",
         0.4, 0.2, 12.0, 0.7, size=25, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.5, 0.04, ACCENT_BLUE)

add_bullet_box(
    sl,
    "Como funciona",
    [
        "Recebe o throughput medido em kbps",
        "Calcula safety_limit = throughput × 0,85  (margem de 15%)",
        "Ordena qualidades por bitrate (crescente)",
        "Escolhe a maior qualidade com bitrate ≤ safety_limit",
        "Fallback: 240p (mínimo) se nenhuma qualidade couber",
        "Registra cada decisão com timestamp no histórico",
    ],
    0.4, 1.1, 5.8, 4.0,
    title_color=ACCENT_BLUE, bg_color=BG_CARD
)

add_code_block(sl, [
    "class RateBasedABR:",
    "    SAFETY_FACTOR = 0.85   # margem de segurança",
    "",
    "    def select_quality(self, throughput_kbps, qualities):",
    "        safety_limit = throughput_kbps * self.SAFETY_FACTOR",
    "",
    "        sorted_q = sorted(qualities, key=lambda q: q['bitrate'])",
    "",
    "        selected = None",
    "        for q in sorted_q:",
    "            if q['bitrate'] <= safety_limit:",
    "                selected = q['name']   # sempre atualiza → pega o maior",
    "",
    "        if selected is None:",
    "            selected = '240p'   # fallback mínimo",
    "",
    "        self._record_decision(throughput_kbps, safety_limit, selected)",
    "        return selected",
], 6.4, 1.1, 6.5, 5.2, font_size=10)

add_bullet_box(
    sl,
    "Utilitários do histórico",
    [
        "get_decision_history() — todas as decisões",
        "get_quality_switches() — apenas as mudanças de qualidade",
        "get_decision_count(quality) — contagem por qualidade",
    ],
    0.4, 5.25, 5.8, 2.0,
    title_color=ACCENT_GREEN, bg_color=BG_CARD, bullet_size=13
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – BufferManager
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_PURP)
slide_number_label(sl, 7)

add_text(sl, "BufferManager — Gestão do Buffer de Playback",
         0.4, 0.2, 12.0, 0.7, size=25, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.5, 0.04, ACCENT_PURP)

add_bullet_box(
    sl,
    "Constantes e estado",
    [
        "max_buffer = 60 s (padrão configurável)",
        "MIN_BUFFER_TO_PLAY = 2 s → mínimo para iniciar",
        "REBUFFER_THRESHOLD = 0 s → dispara rebuffering",
        "current_buffer — nível atual em segundos",
        "rebuffer_count / rebuffer_history — contagem e log",
    ],
    0.4, 1.1, 5.8, 3.4,
    title_color=ACCENT_PURP, bg_color=BG_CARD
)

add_code_block(sl, [
    "def add_segment(self, duration):",
    "    # Adiciona chunk ao buffer (capped em max_buffer)",
    "    self.current_buffer = min(",
    "        self.current_buffer + duration, self.max_buffer",
    "    )",
    "",
    "def consume(self, time_elapsed):",
    "    old = self.current_buffer",
    "    self.current_buffer = max(0.0, old - time_elapsed)",
    "    self.total_consumed += time_elapsed",
    "    # Detecta rebuffering: buffer vai a zero",
    "    if old > 0 and self.current_buffer <= 0:",
    "        self._record_rebuffer()",
    "",
    "def can_play(self):",
    "    return self.current_buffer >= self.MIN_BUFFER_TO_PLAY",
], 6.4, 1.1, 6.5, 4.8, font_size=10.5)

add_bullet_box(
    sl,
    "Ciclo de vida do buffer",
    [
        "add_segment() → segmento chega da rede",
        "consume() → player consome em tempo real",
        "is_rebuffering() → buffer zerou → player para",
        "can_play() → buffer ≥ 2 s → retoma",
        "get_stats() → snapshot completo de métricas",
    ],
    0.4, 4.65, 5.8, 2.6,
    title_color=ACCENT_ORANGE, bg_color=BG_CARD, bullet_size=13
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – ThroughputMeter
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_GREEN)
slide_number_label(sl, 8)

add_text(sl, "ThroughputMeter — Medição de Banda",
         0.4, 0.2, 12.0, 0.7, size=26, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 4.5, 0.04, ACCENT_GREEN)

add_code_block(sl, [
    "def start_measurement(self):",
    "    self.current_start_time = time.time()",
    "",
    "def stop_measurement(self, bytes_downloaded):",
    "    elapsed = time.time() - self.current_start_time",
    "    if elapsed <= 0:",
    "        elapsed = 0.001   # evita divisão por zero",
    "",
    "    # bytes → bits → kbits / tempo",
    "    throughput_kbps = (bytes_downloaded * 8) / elapsed / 1000",
    "",
    "    measurement = ThroughputMeasurement(",
    "        bytes_downloaded=bytes_downloaded,",
    "        time_elapsed=elapsed,",
    "        throughput_kbps=throughput_kbps,",
    "        timestamp=datetime.now(),",
    "    )",
    "    self.history.append(measurement)   # deque(maxlen=5)",
    "    return measurement",
], 0.4, 1.1, 7.0, 5.7, font_size=10.5)

add_bullet_box(
    sl,
    "Análises disponíveis",
    [
        "get_average_throughput() — média das últimas N medições",
        "get_min / get_max — mínimo e máximo no histórico",
        "get_jitter() — desvio padrão (variabilidade da rede)",
        "get_throughput_trend() — increasing / decreasing / stable",
        "history_size padrão = 5 medições (sliding window)",
    ],
    7.6, 1.1, 5.3, 3.6,
    title_color=ACCENT_GREEN, bg_color=BG_CARD
)

add_bullet_box(
    sl,
    "Fórmula de throughput",
    [
        "throughput_kbps = (bytes × 8) ÷ tempo_s ÷ 1000",
        "Janela deslizante de 5 medições (deque com maxlen)",
        "Jitter = desvio padrão das medições na janela",
    ],
    7.6, 4.9, 5.3, 2.4,
    title_color=ACCENT_ORANGE, bg_color=BG_CARD, bullet_size=13
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – MetricsRecorder
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_ORANGE)
slide_number_label(sl, 9)

add_text(sl, "MetricsRecorder — Coleta e Persistência de Dados",
         0.4, 0.2, 12.0, 0.7, size=25, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.5, 0.04, ACCENT_ORANGE)

add_code_block(sl, [
    "class MetricsRecorder:",
    "    headers = [",
    "        'segment_id', 'timestamp', 'quality',",
    "        'bitrate_kbps', 'throughput_kbps',",
    "        'buffer_level_secs', 'rebuffering_occurred', 'quality_changed'",
    "    ]",
    "",
    "    def record_segment(self, segment_data):",
    "        self.buffer_data.append(segment_data)",
    "        if len(self.buffer_data) >= self.batch_size:",
    "            self.flush()   # grava em lote",
    "",
    "    def flush(self):",
    "        with open(self.filepath, 'a') as f:",
    "            writer = csv.DictWriter(f, fieldnames=self.headers)",
    "            writer.writerows(self.buffer_data)",
    "        self.buffer_data.clear()",
    "",
    "    def close(self):",
    "        self.flush()   # garante dados residuais",
], 0.4, 1.1, 7.0, 5.9, font_size=10.5)

add_bullet_box(
    sl,
    "Design de desempenho",
    [
        "Escrita em lote (batch_size = 5 padrão)",
        "DictWriter mapeia chaves → evita bugs de ordem",
        "Context manager (__enter__ / __exit__) → close() automático",
        "Nome do arquivo com timestamp → nunca sobrescreve logs",
        "Cria diretório de saída automaticamente se não existir",
    ],
    7.6, 1.1, 5.3, 3.6,
    title_color=ACCENT_ORANGE, bg_color=BG_CARD
)

add_bullet_box(
    sl,
    "8 colunas do CSV",
    [
        "segment_id, timestamp, quality, bitrate_kbps",
        "throughput_kbps, buffer_level_secs",
        "rebuffering_occurred, quality_changed",
    ],
    7.6, 4.9, 5.3, 2.4,
    title_color=ACCENT_BLUE, bg_color=BG_CARD, bullet_size=13
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – run_simulation.py
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_ORANGE)
slide_number_label(sl, 10)

add_text(sl, "run_simulation.py — Simulação Integrada",
         0.4, 0.2, 12.0, 0.7, size=26, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.0, 0.04, ACCENT_ORANGE)

add_code_block(sl, [
    "# Perfil de rede que simula oscilações reais",
    "network_profile = (",
    "    [3000.0] * 5 +   # 3 Mbps  — Alta qualidade",
    "    [1500.0] * 5 +   # 1.5 Mbps — Média",
    "    [300.0]  * 6 +   # 300 kbps — Crise / rebuffering",
    "    [900.0]  * 4 +   # 900 kbps — Recuperando",
    "    [4000.0] * 10    # 4 Mbps  — Excelente",
    ")",
    "",
    "for seg_id in range(1, total_segments + 1):",
    "    # Ruído ±15% sobre o perfil",
    "    throughput = network_profile[seg_id-1] * random.uniform(0.85, 1.15)",
    "    quality = abr.select_quality(throughput, qualities)",
    "    buffer_manager.add_segment(segment_duration)",
    "    download_time = (bitrate / throughput) * segment_duration",
    "    buffer_manager.consume(download_time)",
    "    recorder.record_segment({...})",
], 0.4, 1.1, 7.6, 5.7, font_size=10.5)

add_bullet_box(
    sl,
    "Componentes integrados",
    [
        "RateBasedABR → decide qualidade a cada segmento",
        "BufferManager → simula preenchimento e consumo",
        "MetricsRecorder → persiste 8 colunas em CSV",
        "30 segmentos × 4 s = 120 s de vídeo simulado",
        "Saída: logs/metrics_<timestamp>.csv",
    ],
    8.2, 1.1, 4.7, 3.8,
    title_color=ACCENT_ORANGE, bg_color=BG_CARD
)

add_bullet_box(
    sl,
    "Cenário de rede",
    [
        "Fase 1 — 3 Mbps → qualidade 1080p",
        "Fase 2 — 1,5 Mbps → 720p",
        "Fase 3 — 300 kbps → 240p (rebuffering!)",
        "Fase 4 — 900 kbps → recuperação gradual",
        "Fase 5 — 4 Mbps → 1080p estável",
    ],
    8.2, 5.1, 4.7, 2.2,
    title_color=ACCENT_BLUE, bg_color=BG_CARD, bullet_size=13
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Gráficos e Resultados
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_BLUE)
slide_number_label(sl, 11)

add_text(sl, "Visualização — Gráficos Gerados por graphs.py",
         0.4, 0.2, 12.0, 0.7, size=26, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.5, 0.04, ACCENT_BLUE)

graphs = [
    ("throughput_timeline.png", "Throughput vs Qualidade",
     "Linha de throughput (kbps) com barras de qualidade selecionada no eixo duplo.",
     ACCENT_BLUE),
    ("quality_timeline.png",    "Timeline de Qualidade",
     "Gráfico em degrau mostrando cada mudança de qualidade ao longo do tempo.",
     ACCENT_GREEN),
    ("buffer_level.png",        "Nível do Buffer",
     "Evolução do buffer em segundos — identifica períodos de rebuffering.",
     ACCENT_PURP),
    ("quality_distribution.png","Distribuição de Qualidade",
     "Tempo total gasto em cada nível de qualidade durante a sessão.",
     ACCENT_GREEN),
    ("throughput_histogram.png", "Histograma de Throughput",
     "Distribuição de frequência das medições de banda registradas.",
     ACCENT_ORANGE),
]

for i, (fname, title, desc, color) in enumerate(graphs):
    col = i % 3
    row = i // 3
    bx = 0.4 + col * 4.3
    by = 1.2 + row * 3.0

    add_rect(sl, bx, by, 4.0, 2.6, BG_CARD)
    add_rect(sl, bx, by, 4.0, 0.35, color)
    add_text(sl, title, bx + 0.1, by + 0.03, 3.8, 0.28,
             size=13, bold=True, color=WHITE)
    add_text(sl, fname, bx + 0.1, by + 0.42, 3.8, 0.28,
             size=10, color=LIGHT_GRAY, italic=True)
    add_text(sl, desc, bx + 0.1, by + 0.78, 3.8, 1.7,
             size=12, color=LIGHT_GRAY)

add_bullet_box(
    sl,
    "Como gerar os gráficos",
    [
        "python run_simulation.py   →  gera logs/metrics_<ts>.csv",
        "python graphs.py logs/metrics_<ts>.csv --output-dir graphs/",
        "Aceita aliases de colunas — compatível com outros CSVs",
    ],
    0.4, 6.5, 12.5, 0.85,
    title_color=ACCENT_ORANGE, bg_color=BG_CARD, bullet_size=13, title_size=14
)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Testes
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_GREEN)
slide_number_label(sl, 12)

add_text(sl, "Testes Unitários — Cobertura por Módulo",
         0.4, 0.2, 12.0, 0.7, size=26, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 5.0, 0.04, ACCENT_GREEN)

test_modules = [
    ("test_abr.py",             "RateBasedABR",   ACCENT_BLUE,
     ["select_quality com vários throughputs",
      "Fallback para 240p quando banda é insuficiente",
      "Safety factor de 85%",
      "Histórico e switches de qualidade",
      "TypeError / ValueError em entradas inválidas"]),
    ("test_buffer_manager.py",  "BufferManager",  ACCENT_PURP,
     ["add_segment e consume",
      "Detecção de rebuffering",
      "can_play (limiar de 2 s)",
      "Buffer cap em max_buffer",
      "get_stats() e reset()"]),
    ("test_manifest_parser.py", "ManifestParser", ACCENT_ORANGE,
     ["load_from_dict com manifest válido",
      "KeyError em campos ausentes",
      "get_servers ordena por prioridade",
      "get_bitrate por nome de qualidade",
      "download_manifest (mock HTTP)"]),
    ("test_metrics.py",         "ThroughputMeter",ACCENT_GREEN,
     ["start/stop measurement",
      "Cálculo de throughput (bytes×8/t/1000)",
      "Jitter (desvio padrão)",
      "get_throughput_trend",
      "RuntimeError sem start()"]),
]

for i, (fname, cls, color, cases) in enumerate(test_modules):
    col = i % 2
    row = i // 2
    bx = 0.4 + col * 6.4
    by = 1.2 + row * 3.0

    add_rect(sl, bx, by, 6.1, 2.8, BG_CARD)
    add_rect(sl, bx, by, 6.1, 0.35, color)
    add_text(sl, fname, bx + 0.12, by + 0.04, 5.8, 0.27,
             size=12, bold=True, color=WHITE)
    add_text(sl, cls,   bx + 0.12, by + 0.42, 5.8, 0.32,
             size=15, bold=True, color=LIGHT_GRAY)
    for j, case in enumerate(cases):
        add_text(sl, f"✓ {case}", bx + 0.12, by + 0.82 + j * 0.38, 5.8, 0.36,
                 size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Conclusão
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, BG_DARK)
add_accent_bar(sl, ACCENT_BLUE)
slide_number_label(sl, 13)

add_text(sl, "Conclusão",
         0.4, 0.2, 12.0, 0.7, size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 0.95, 2.5, 0.04, ACCENT_BLUE)

cards = [
    ("ABR Funcional",
     "Algoritmo Rate-Based com safety factor de 85% seleciona qualidade\n"
     "de forma conservadora e responsiva às variações de banda.",
     ACCENT_BLUE, 0.4, 1.2),
    ("Buffer Robusto",
     "BufferManager detecta rebuffering, controla preenchimento e\n"
     "fornece snapshot completo de métricas em get_stats().",
     ACCENT_PURP, 6.8, 1.2),
    ("Métricas Completas",
     "MetricsRecorder persiste 8 colunas por segmento em CSV com\n"
     "escrita em lote e context manager para fechamento seguro.",
     ACCENT_ORANGE, 0.4, 3.6),
    ("Visualização Rica",
     "graphs.py gera 5 gráficos distintos: timeline, distribuição,\n"
     "histograma, buffer e correlação throughput × qualidade.",
     ACCENT_GREEN, 6.8, 3.6),
]

for title, body, color, bx, by in cards:
    add_rect(sl, bx, by, 6.1, 2.1, BG_CARD)
    add_rect(sl, bx, by, 0.12, 2.1, color)
    add_text(sl, title, bx + 0.28, by + 0.18, 5.7, 0.5,
             size=18, bold=True, color=WHITE)
    add_text(sl, body, bx + 0.28, by + 0.75, 5.7, 1.2,
             size=14, color=LIGHT_GRAY)

add_rect(sl, 0.4, 5.85, 12.53, 0.04, ACCENT_BLUE)
add_text(sl,
         "O sistema implementa o ciclo completo de streaming adaptativo:\n"
         "medição de banda → seleção de qualidade → gestão de buffer → coleta de métricas → visualização.",
         0.4, 6.0, 12.53, 1.0,
         size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)


# ── Salvar ────────────────────────────────────────────────────────────────────
output = "/home/gabriel_pc/Testes/TR2/adaptive-http-streaming-main/Adaptive_HTTP_Streaming.pptx"
prs.save(output)
print(f"Apresentação salva em: {output}")
